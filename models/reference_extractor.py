from pptx import Presentation

def extract_reference_structure(file_path: str) -> str:
    prs = Presentation(file_path)
    structure = []

    for slide in prs.slides:
        slide_dict = {"placeholders": []}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            placeholder = shape.name
            slide_dict["placeholders"].append(placeholder)
        structure.append(slide_dict)

    return str(structure)
