from pptx.util import Pt
from pptx import Presentation

class PPTGenerator:
    def __init__(self, template_manager):
        self.template_manager = template_manager
        self.output = Presentation()

    def populate_slide(self, slide_content):
        template = self.template_manager.load_template(slide_content.slide_type)
        slide_layout = template.slide_layouts[0]
        slide = self.output.slides.add_slide(slide_layout)

        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            placeholder_name = shape.name.lower()
            if placeholder_name in slide_content.content:
                shape.text = slide_content.content[placeholder_name]
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
                        run.font.name = "Calibri"

        if slide_content.image_paths:
            for ph_name, img_path in slide_content.image_paths.items():
                for shape in slide.shapes:
                    if ph_name in shape.name.lower():
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        slide.shapes.add_picture(img_path, left, top, width, height)
                        slide.shapes._spTree.remove(shape._element)
                        break

    def save(self, output_path):
        self.output.save(output_path)
